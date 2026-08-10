"""Convertidores por formato de origen → Markdown GFM."""

from __future__ import annotations

import csv
import json
import re
import shutil
from email import policy
from email.parser import BytesParser
from html import unescape
from io import StringIO
from pathlib import Path
from typing import Callable

from detect import DetectedType, UnsupportedFormatError
from markdown_out import ConversionResult, assets_dir_for, gfm_table, heading, relative_asset_path
from ocr_utils import ocr_image_file, ocr_pdf_to_markdown, pdf_looks_like_scan

ProgressCallback = Callable[[str], None]


def convert_detected(
    path: Path,
    detected: DetectedType,
    *,
    md_path: Path,
    ocr_mode: str = "auto",
    progress: ProgressCallback | None = None,
) -> ConversionResult:
    kind = detected.kind
    handlers = {
        "pdf": _convert_pdf,
        "docx": _convert_docx,
        "pptx": _convert_pptx,
        "xlsx": _convert_xlsx,
        "odt": _convert_odt,
        "rtf": _convert_rtf,
        "txt": _convert_txt,
        "csv": _convert_csv,
        "tsv": _convert_tsv,
        "json": _convert_json,
        "html": _convert_html_file,
        "epub": _convert_epub,
        "eml": _convert_eml,
        "msg": _convert_msg,
        "png": _convert_image,
        "jpg": _convert_image,
        "tiff": _convert_image,
    }
    handler = handlers.get(kind)
    if not handler:
        raise UnsupportedFormatError(f"No hay convertidor para el tipo «{kind}».")
    result = handler(path, md_path=md_path, ocr_mode=ocr_mode, progress=progress)
    result.source_format = kind
    result.source_name = path.name
    if not result.title:
        result.title = path.stem
    return result


def convert_url(url: str, *, md_path: Path, progress: ProgressCallback | None = None) -> ConversionResult:
    import urllib.request

    if progress:
        progress(f"Descargando URL…")
    req = urllib.request.Request(url, headers={"User-Agent": "file2md/2.2"})
    with urllib.request.urlopen(req, timeout=60) as resp:  # noqa: S310 - URL pedida por el usuario
        raw = resp.read()
        content_type = resp.headers.get("Content-Type", "text/html")
    text = raw.decode("utf-8", errors="replace")
    if "html" not in content_type.lower() and not re.search(r"<html", text[:1000], re.I):
        raise UnsupportedFormatError(
            f"La URL no devolvió HTML convertible (Content-Type: {content_type})."
        )
    body, title = _html_to_markdown(text)
    return ConversionResult(
        markdown_body=body,
        title=title or url,
        source_format="url",
        source_name=url,
        extra_meta={"source_url": url},
    )


def _report(progress: ProgressCallback | None, msg: str) -> None:
    if progress:
        progress(msg)


# ----- PDF -----


def _convert_pdf(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    import pymupdf
    import pymupdf4llm

    _report(progress, "Analizando PDF…")
    use_ocr = ocr_mode == "always" or (ocr_mode == "auto" and pdf_looks_like_scan(path))
    assets = assets_dir_for(md_path)
    image_paths: list[Path] = []

    if use_ocr:
        _report(progress, "PDF sin capa de texto. Aplicando OCR…")
        body = ocr_pdf_to_markdown(path, progress=progress)
    else:
        _report(progress, "Extrayendo texto y estructura del PDF…")
        # write_images ayuda a preservar figuras; orden de lectura vía layout de pymupdf4llm.
        assets.mkdir(parents=True, exist_ok=True)
        try:
            body = pymupdf4llm.to_markdown(
                str(path),
                write_images=True,
                image_path=str(assets),
                image_format="png",
            )
        except TypeError:
            body = pymupdf4llm.to_markdown(str(path))
        if not str(body).strip() and ocr_mode == "auto":
            _report(progress, "Sin texto útil. Reintentando con OCR…")
            body = ocr_pdf_to_markdown(path, progress=progress)
        else:
            # Reescribir rutas de imagen a relativas desde el .md
            body = _normalize_asset_links(str(body), md_path, assets)
            if assets.exists():
                image_paths = sorted(assets.glob("**/*.*"))

    title = path.stem
    author = None
    with pymupdf.open(path) as doc:
        meta = doc.metadata or {}
        title = meta.get("title") or title
        author = meta.get("author") or None

    return ConversionResult(
        markdown_body=str(body).strip(),
        title=title,
        author=author,
        source_format="pdf",
        asset_files=image_paths,
    )


def _normalize_asset_links(markdown: str, md_path: Path, assets: Path) -> str:
    if not assets.exists():
        return markdown

    def repl(match: re.Match[str]) -> str:
        alt, src = match.group(1), match.group(2)
        src_path = Path(src)
        if not src_path.is_file():
            candidate = assets / src_path.name
            if candidate.is_file():
                src_path = candidate
            else:
                return match.group(0)
        rel = relative_asset_path(md_path, src_path)
        return f"![{alt}]({rel})"

    return re.sub(r"!\[([^\]]*)\]\(([^)]+)\)", repl, markdown)


# ----- DOCX -----


def _convert_docx(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    from docx import Document
    from docx.document import Document as DocumentObject
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    _report(progress, "Convirtiendo DOCX…")
    doc = Document(str(path))
    assets = assets_dir_for(md_path)
    assets.mkdir(parents=True, exist_ok=True)
    parts: list[str] = []
    image_paths: list[Path] = []
    img_i = 0
    rels = doc.part.rels

    def style_level(style_name: str) -> int | None:
        name = (style_name or "").lower()
        for n in range(1, 7):
            if f"heading {n}" in name or f"título {n}" in name or name.startswith(f"heading{n}"):
                return n
        return None

    def extract_images_from_paragraph(paragraph: Paragraph) -> list[str]:
        nonlocal img_i
        links: list[str] = []
        for run in paragraph.runs:
            for blip in run._element.findall(".//" + qn("a:blip")):
                embed = blip.get(qn("r:embed"))
                if not embed or embed not in rels:
                    continue
                img_i += 1
                rel = rels[embed]
                blob = rel.target_part.blob
                ext = Path(rel.target_ref).suffix or ".png"
                out_img = assets / f"image_{img_i:03d}{ext}"
                out_img.write_bytes(blob)
                image_paths.append(out_img)
                links.append(f"![]({relative_asset_path(md_path, out_img)})")
        return links

    def iter_block_items(parent: DocumentObject):
        body = parent.element.body
        for child in body.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, parent)
            elif child.tag == qn("w:tbl"):
                yield Table(child, parent)

    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            text = (block.text or "").strip()
            images = extract_images_from_paragraph(block)
            level = style_level(block.style.name if block.style else "")
            style_name = (block.style.name if block.style else "") or ""
            if level and text:
                parts.append(heading(level, text))
            elif "list" in style_name.lower() and text:
                parts.append(f"- {text}")
            elif text:
                # Enlaces: python-docx no siempre los expone; conservar texto.
                parts.append(text)
            parts.extend(images)
        elif isinstance(block, Table):
            rows_data: list[list[str]] = []
            for row in block.rows:
                rows_data.append([cell.text.strip() for cell in row.cells])
            if rows_data:
                parts.append(gfm_table(rows_data[0], rows_data[1:]))

    core = doc.core_properties
    return ConversionResult(
        markdown_body="\n\n".join(parts).strip(),
        title=core.title or path.stem,
        author=core.author or None,
        date=core.created.isoformat() if core.created else None,
        asset_files=image_paths,
    )


# ----- PPTX -----


def _convert_pptx(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    _report(progress, "Convirtiendo PPTX…")
    prs = Presentation(str(path))
    assets = assets_dir_for(md_path)
    assets.mkdir(parents=True, exist_ok=True)
    parts: list[str] = []
    image_paths: list[Path] = []
    img_i = 0

    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(heading(2, f"Diapositiva {idx}"))
        # Orden aproximado de lectura: top-to-bottom, left-to-right
        shapes = sorted(
            slide.shapes,
            key=lambda s: (round(getattr(s, "top", 0) or 0), round(getattr(s, "left", 0) or 0)),
        )
        for shape in shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if not text:
                        text = (para.text or "").strip()
                    if not text:
                        continue
                    level = (para.level or 0) + 1
                    if level > 1:
                        parts.append(f"{'  ' * (level - 1)}- {text}")
                    else:
                        parts.append(text)
            if shape.has_table:
                table = shape.table
                rows = []
                for r in table.rows:
                    rows.append([c.text.strip() for c in r.cells])
                if rows:
                    parts.append(gfm_table(rows[0], rows[1:]))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_i += 1
                image = shape.image
                ext = image.ext or "png"
                out_img = assets / f"slide{idx}_image_{img_i:03d}.{ext}"
                out_img.write_bytes(image.blob)
                image_paths.append(out_img)
                parts.append(f"![]({relative_asset_path(md_path, out_img)})")

    return ConversionResult(
        markdown_body="\n\n".join(parts).strip(),
        title=path.stem,
        asset_files=image_paths,
    )


# ----- XLSX -----


def _convert_xlsx(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    from openpyxl import load_workbook

    _report(progress, "Convirtiendo XLSX…")
    wb = load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(heading(2, sheet.title))
        rows = []
        for row in sheet.iter_rows(values_only=True):
            if row is None or all(c is None or str(c).strip() == "" for c in row):
                continue
            rows.append(["" if c is None else str(c) for c in row])
        if rows:
            parts.append(gfm_table(rows[0], rows[1:]))
    wb.close()
    return ConversionResult(markdown_body="\n\n".join(parts).strip(), title=path.stem)


# ----- ODT -----


def _convert_odt(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    from odf import teletype, text
    from odf.opendocument import load

    _report(progress, "Convirtiendo ODT…")
    doc = load(str(path))
    parts: list[str] = []
    for el in doc.getElementsByType(text.H):
        level = int(el.getAttribute("outlinelevel") or "1")
        parts.append(heading(level, teletype.extractText(el)))
    # Párrafos (incluye algunos ya vistos en headings; filtramos vacíos)
    for el in doc.getElementsByType(text.P):
        t = teletype.extractText(el).strip()
        if t:
            parts.append(t)
    # Mejor aproximación: recorrer body en orden
    body_parts: list[str] = []
    body = doc.body
    for child in body.childNodes:
        name = getattr(child, "tagName", "") or ""
        if name.endswith(":h") or name == "text:h":
            level = int(child.getAttribute("outlinelevel") or "1")
            body_parts.append(heading(level, teletype.extractText(child)))
        elif name.endswith(":p") or name == "text:p":
            t = teletype.extractText(child).strip()
            if t:
                body_parts.append(t)
        elif name.endswith(":list") or name == "text:list":
            for item in child.getElementsByType(text.ListItem):
                t = teletype.extractText(item).strip()
                if t:
                    body_parts.append(f"- {t}")
        elif "table" in name:
            rows = []
            for row_el in child.getElementsByType(text.P):
                pass
            # fallback simple de tabla via XML rows
            for row_el in child.childNodes:
                rname = getattr(row_el, "tagName", "") or ""
                if "table-row" in rname:
                    cells = []
                    for cell in row_el.childNodes:
                        cname = getattr(cell, "tagName", "") or ""
                        if "table-cell" in cname:
                            cells.append(teletype.extractText(cell).strip())
                    if cells:
                        rows.append(cells)
            if rows:
                body_parts.append(gfm_table(rows[0], rows[1:]))

    content = "\n\n".join(body_parts if body_parts else parts).strip()
    return ConversionResult(markdown_body=content, title=path.stem)


# ----- RTF / TXT / CSV / TSV / JSON -----


def _convert_rtf(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    from striprtf.striprtf import rtf_to_text

    _report(progress, "Convirtiendo RTF…")
    raw = path.read_text(encoding="utf-8", errors="ignore")
    text = rtf_to_text(raw).strip()
    return ConversionResult(markdown_body=text, title=path.stem)


def _convert_txt(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    _report(progress, "Leyendo texto…")
    text = path.read_text(encoding="utf-8-sig", errors="replace")
    return ConversionResult(markdown_body=text.strip(), title=path.stem)


def _convert_csv(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    _report(progress, "Convirtiendo CSV…")
    text = path.read_text(encoding="utf-8-sig", errors="replace")
    reader = csv.reader(StringIO(text))
    rows = [list(r) for r in reader if any(c.strip() for c in r)]
    if not rows:
        raise RuntimeError("El CSV no tiene filas con datos.")
    return ConversionResult(markdown_body=gfm_table(rows[0], rows[1:]), title=path.stem)


def _convert_tsv(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    _report(progress, "Convirtiendo TSV…")
    text = path.read_text(encoding="utf-8-sig", errors="replace")
    reader = csv.reader(StringIO(text), delimiter="\t")
    rows = [list(r) for r in reader if any(c.strip() for c in r)]
    if not rows:
        raise RuntimeError("El TSV no tiene filas con datos.")
    return ConversionResult(markdown_body=gfm_table(rows[0], rows[1:]), title=path.stem)


def _convert_json(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    _report(progress, "Convirtiendo JSON…")
    data = json.loads(path.read_text(encoding="utf-8-sig"))
    pretty = json.dumps(data, ensure_ascii=False, indent=2)
    body = f"```json\n{pretty}\n```"
    return ConversionResult(markdown_body=body, title=path.stem)


# ----- HTML -----


def _convert_html_file(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    _report(progress, "Convirtiendo HTML…")
    html = path.read_text(encoding="utf-8-sig", errors="replace")
    body, title = _html_to_markdown(html)
    return ConversionResult(markdown_body=body, title=title or path.stem)


def _html_to_markdown(html: str) -> tuple[str, str | None]:
    try:
        from markdownify import markdownify as md
    except ImportError:
        md = None

    title = None
    m = re.search(r"<title[^>]*>(.*?)</title>", html, flags=re.I | re.S)
    if m:
        title = unescape(re.sub(r"\s+", " ", m.group(1))).strip()

    if md:
        body = md(
            html,
            heading_style="ATX",
            bullets="-",
            code_language="",
            strip=["script", "style"],
        )
        return body.strip(), title

    # Fallback mínimo sin markdownify
    text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", "", html)
    text = re.sub(r"(?i)<br\s*/?>", "\n", text)
    text = re.sub(r"(?i)</p>", "\n\n", text)
    text = re.sub(r"(?i)<h([1-6])[^>]*>", lambda m: "\n" + "#" * int(m.group(1)) + " ", text)
    text = re.sub(r"(?i)</h[1-6]>", "\n\n", text)
    text = re.sub(r"(?s)<[^>]+>", "", text)
    return unescape(text).strip(), title


# ----- EPUB -----


def _convert_epub(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    from ebooklib import ITEM_DOCUMENT, epub

    _report(progress, "Convirtiendo EPUB…")
    book = epub.read_epub(str(path))
    title = path.stem
    author = None
    if book.get_metadata("DC", "title"):
        title = book.get_metadata("DC", "title")[0][0]
    if book.get_metadata("DC", "creator"):
        author = book.get_metadata("DC", "creator")[0][0]

    parts: list[str] = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        html = item.get_content().decode("utf-8", errors="replace")
        body, _ = _html_to_markdown(html)
        if body.strip():
            parts.append(body.strip())

    return ConversionResult(
        markdown_body="\n\n".join(parts).strip(),
        title=title,
        author=author,
    )


# ----- Email -----


def _convert_eml(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    _report(progress, "Convirtiendo EML…")
    msg = BytesParser(policy=policy.default).parsebytes(path.read_bytes())
    subject = msg.get("subject", path.stem)
    author = msg.get("from")
    date = msg.get("date")
    parts = [
        heading(1, subject or path.stem),
        f"**De:** {author}" if author else "",
        f"**Para:** {msg.get('to')}" if msg.get("to") else "",
        f"**Fecha:** {date}" if date else "",
    ]
    body = msg.get_body(preferencelist=("plain", "html"))
    if body is not None:
        content = body.get_content()
        if body.get_content_type() == "text/html":
            md_body, _ = _html_to_markdown(content)
            parts.append(md_body)
        else:
            parts.append(content.strip())
    else:
        parts.append("(Sin cuerpo de mensaje)")
    return ConversionResult(
        markdown_body="\n\n".join(p for p in parts if p).strip(),
        title=subject,
        author=author,
        date=str(date) if date else None,
    )


def _convert_msg(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    import extract_msg

    _report(progress, "Convirtiendo MSG…")
    msg = extract_msg.Message(str(path))
    try:
        subject = msg.subject or path.stem
        author = msg.sender
        date = str(msg.date) if msg.date else None
        body = (msg.body or "").strip()
        if not body and msg.htmlBody:
            body, _ = _html_to_markdown(
                msg.htmlBody.decode("utf-8", errors="replace")
                if isinstance(msg.htmlBody, bytes)
                else str(msg.htmlBody)
            )
        parts = [
            heading(1, subject),
            f"**De:** {author}" if author else "",
            f"**Fecha:** {date}" if date else "",
            body or "(Sin cuerpo de mensaje)",
        ]
        return ConversionResult(
            markdown_body="\n\n".join(p for p in parts if p).strip(),
            title=subject,
            author=author,
            date=date,
        )
    finally:
        msg.close()


# ----- Images -----


def _convert_image(
    path: Path,
    *,
    md_path: Path,
    ocr_mode: str,
    progress: ProgressCallback | None,
) -> ConversionResult:
    _report(progress, "Imagen detectada. Aplicando OCR…")
    assets = assets_dir_for(md_path)
    assets.mkdir(parents=True, exist_ok=True)
    dest = assets / path.name
    shutil.copy2(path, dest)
    text = ocr_image_file(path, progress=progress)
    body = f"![]({relative_asset_path(md_path, dest)})\n\n{text}".strip()
    return ConversionResult(
        markdown_body=body,
        title=path.stem,
        asset_files=[dest],
        extra_meta={"ocr": True},
    )
